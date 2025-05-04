"""
PPT处理模块

提供PPT文件的处理功能，将PPT文件转换为结构化数据
"""

import logging
import os
from typing import Dict, Any
import pptx
from pptx import Presentation

# 配置日志
logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)


def process_ppt(ppt_path: str) -> Dict[str, Any]:
    """
    处理PPT文件并转换为结构化数据
    
    参数:
        ppt_path (str): PPT文件路径
        
    返回:
        Dict[str, Any]: 处理结果，包含processed_data和processed_content
    """
    logger.info(f"开始处理PPT文件: {ppt_path}")
    
    try:
        # 检查文件是否存在
        if not os.path.exists(ppt_path):
            logger.error(f"PPT文件不存在: {ppt_path}")
            return {"error": f"PPT文件不存在: {ppt_path}"}
        
        # 打开PPT文件
        presentation = Presentation(ppt_path)
        
        # 提取PPT内容
        slides_content = []
        for i, slide in enumerate(presentation.slides):
            slide_content = {
                "slide_number": i + 1,
                "title": slide.shapes.title.text if slide.shapes.title else f"幻灯片 {i+1}",
                "content": []
            }
            
            # 提取幻灯片中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content["content"].append(shape.text)
            
            slides_content.append(slide_content)
        
        # 生成处理后的内容
        processed_content = ""
        for slide in slides_content:
            # 添加幻灯片标题作为一级标题
            processed_content += f"# {slide['title']}\n\n"
            
            # 处理幻灯片内容，确保每个内容块之间有适当的分隔
            for content_block in slide['content']:
                # 移除内容前后的空白字符
                content_block = content_block.strip()
                if content_block:
                    processed_content += f"{content_block}\n\n"
            
            # 在幻灯片之间添加额外的分隔
            processed_content += "\n"
        
        # 返回处理结果
        result = {
            "type": "ppt",
            "path": ppt_path,
            "processed_data": slides_content,
            "processed_content": processed_content,
            "collection_id": os.path.basename(ppt_path)
        }
        
        logger.info("PPT处理成功完成")
        return result
        
    except Exception as e:
        logger.error(f"处理PPT文件时出错: {e}")
        return {
            "type": "ppt",
            "path": ppt_path,
            "error": f"处理PPT文件时出错: {str(e)}",
            "processed_content": f"[PPT内容处理失败: {ppt_path}]",
            "collection_id": os.path.basename(ppt_path)
        }